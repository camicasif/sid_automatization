import glob

from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from PIL import Image
import os
import shutil
from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import copy
file_name = 'fotos_antenas'
OFFSET_BUSQUEDA = 12  # Ajusta según necesidad

def generar_presentacion():
    # Cargar plantilla
    prs = Presentation("Plantilla.pptx")

    # Recorrer carpetas principales
    for carpeta in os.listdir(file_name):
        carpeta_path = os.path.join(file_name, carpeta)
        if not os.path.isdir(carpeta_path):
            continue

        # 1. Añadir diapositiva de título con nombre de carpeta
        slide_titulo = prs.slides.add_slide(prs.slide_layouts[0])
        titulo = slide_titulo.shapes.title
        titulo.text = carpeta

        # 2. Procesar antenas en esta carpeta
        antenas = sorted(
            [d for d in os.listdir(carpeta_path)
             if d.startswith('Antena_') and os.path.isdir(os.path.join(carpeta_path, d))],
            key=lambda x: int(x.split('_')[1])
        )

        for antena in antenas:
            antena_path = os.path.join(carpeta_path, antena)

            # 3. Clonar slide de plantilla para esta antena
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Usar segundo layout de plantilla

            # 4. Insertar imágenes
            img_left = Inches(1)
            img_top = Inches(2)
            img_width = Inches(3)
            img_height = Inches(2)
            current_x = img_left

            for img_file in sorted(os.listdir(antena_path)):
                if img_file.endswith('.png'):
                    img_path = os.path.join(antena_path, img_file)

                    # Añadir imagen
                    pic = slide.shapes.add_picture(
                        img_path,
                        current_x, img_top,
                        width=img_width,
                        height=img_height
                    )

                    # Añadir texto de tecnologías
                    tech_text = obtener_tecnologias(img_file)
                    textbox = slide.shapes.add_textbox(
                        current_x, img_top + img_height + Inches(0.3),
                        img_width, Inches(0.5)
                    )
                    text_frame = textbox.text_frame
                    p = text_frame.paragraphs[0]
                    p.text = tech_text
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

                    current_x += img_width + Inches(0.5)

    prs.save('presentaciones.pptx')

def obtener_tecnologias(nombre_archivo):
    if '(' in nombre_archivo and ')' in nombre_archivo:
        tech_part = nombre_archivo.split('(')[-1].split(')')[0]
        return tech_part.replace('-', ' + ')
    return ""

def buscar_antenas_por_sectores(excel_path, lista_sectores, lista_antenas, output_folder):
    wb = load_workbook(excel_path, data_only=True)
    sheet = wb.worksheets[7]

    imagenes_dict = {}
    for img in sheet._images:
        pos = img.anchor._from
        excel_row = pos.row + 1
        excel_col = pos.col + 1
        imagenes_dict[(excel_row, excel_col)] = img

    merged_ranges = list(sheet.merged_cells.ranges)

    # Crear carpetas Antena_X dentro del proyecto
    for antena in lista_antenas:
        folder_path = os.path.join(output_folder, f"Antena_{antena}")
        os.makedirs(folder_path, exist_ok=True)
    # Buscar todas las combinaciones
    for sector in lista_sectores:
        for antena in lista_antenas:
            try:
                frase_busqueda = f"foto general de la antena {antena} sector {sector}"
                print(f"\nBuscando: {frase_busqueda}")

                # Buscar celda con texto optimizado (Recomendación 2)
                target_cell = None
                descripcion_tecnica = None

                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and frase_busqueda in str(cell.value).lower():
                            target_cell = cell
                            break
                    if target_cell:
                        celda_encontrada = f"{get_column_letter(target_cell.column)}{target_cell.row}"
                        print(f"Texto encontrado en la celda: {celda_encontrada}")

                        # Extraer texto después de los dos puntos
                        texto_completo = str(target_cell.value)
                        if ":" in texto_completo:
                            _, descripcion = texto_completo.split(":", 1)
                            descripcion_tecnica = descripcion.strip()[:30]  # Limitar a 30 caracteres
                            descripcion_tecnica = descripcion_tecnica.replace("/", "-").replace("\\", "-")

                        else:
                            print("No se encontró el caracter ':' en el texto")

                        break

                if not target_cell:
                    print(f"No encontrado: {frase_busqueda}")
                    continue

                celda_encontrada = f"{get_column_letter(target_cell.column)}{target_cell.row}"
                print(f"Texto encontrado en: {celda_encontrada}")

                # Detectar celdas combinadas (optimizado con cache)
                merged_range = None
                for merged in merged_ranges:
                    if (merged.min_row <= target_cell.row <= merged.max_row and
                            merged.min_col <= target_cell.column <= merged.max_col):
                        merged_range = merged
                        break

                # Definir rango de búsqueda (Recomendación 5)
                rango_filas = range(max(1, target_cell.row - OFFSET_BUSQUEDA), target_cell.row)
                start_col = merged_range.min_col if merged_range else target_cell.column
                end_col = merged_range.max_col if merged_range else target_cell.column
                rango_columnas = range(start_col, end_col + 1)

                # Mostrar rango de búsqueda
                print(f"Buscando en: columnas {get_column_letter(start_col)}-{get_column_letter(end_col)}, "
                      f"filas {rango_filas.start}-{rango_filas.stop-1}")

                # Buscar en el diccionario de imágenes (Optimización clave)
                imagen_encontrada = False
                for fila in rango_filas:
                    for col in rango_columnas:
                        if (fila, col) in imagenes_dict:
                            img = imagenes_dict[(fila, col)]
                            folder = os.path.join(output_folder, f"Antena_{antena}")
                            if descripcion_tecnica:
                                filename = f"Antena_{antena}_Sector_{sector}_({descripcion_tecnica}).png"
                            else:
                                filename = f"Antena_{antena}_Sector_{sector}.png"
                            output_path = os.path.join(folder, filename)

                            try:
                                img_data = img._data()
                                with open(output_path, "wb") as f:
                                    f.write(img_data)

                                with Image.open(output_path) as img_pil:
                                    img_pil.verify()

                                print(f"Imagen guardada en: {output_path}")
                                imagen_encontrada = True
                                break

                            except Exception as e:
                                print(f"Error guardando imagen: {str(e)}")

                        if imagen_encontrada:
                            break

                if not imagen_encontrada:
                    print(f"¡Imagen no encontrada en el rango especificado!")

            except Exception as e:
                print(f"Error procesando {frase_busqueda}: {str(e)}")
                continue

def procesar_excels():
    lista_sectores = ['a', 'b', 'c']
    lista_antenas = [1, 2, 3, 4]

    for excel_path in glob.glob(os.path.join("TSS", "*.xls*")):
        try:
            print(f"\nProcesando: {excel_path}")

            # Extraer nombre del proyecto (ej: TSS - FUNDO VENTILLAS_ver.B.xslm)
            nombre_base = os.path.splitext(os.path.basename(excel_path))[0]
            if len(nombre_base) > 4:
                nombre_carpeta = nombre_base[4:].split(".")[0]
            else:
                nombre_carpeta = nombre_base

            output_folder = os.path.join(file_name, nombre_carpeta)
            shutil.rmtree(output_folder, ignore_errors=True)
            os.makedirs(output_folder, exist_ok=True)

            buscar_antenas_por_sectores(excel_path, lista_sectores, lista_antenas, output_folder)
            print(f"Imágenes guardadas en: {output_folder}")

        except Exception as e:
            print(f"Error procesando {excel_path}: {str(e)}")

if __name__ == "__main__":
    procesar_excels()
    generar_presentacion()
    print("\nProceso completado. Estructura generada en 'fotos_antenas/'")


